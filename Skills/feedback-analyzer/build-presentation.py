"""
Reference script for building minimalist PowerPoint presentations from feedback data.
Used by the feedback-analyzer skill.

This script demonstrates the visual style, layout, and structure for both
the full themes presentation and the polling questions presentation.

Usage:
  This is a reference template. The agent should adapt it based on the actual
  feedback data and user requirements.

Requirements:
  pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# COLOR PALETTE — minimalist, professional
# ============================================================
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF7, 0xF7, 0xF7)
DARK        = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT   = RGBColor(0x5A, 0x5A, 0x6E)
LIGHT_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY    = RGBColor(0x9A, 0x9A, 0xA8)
ACCENT      = RGBColor(0x2B, 0x57, 0x97)     # Muted blue
ACCENT_LIGHT= RGBColor(0xE8, 0xEF, 0xF7)

# Priority colors
P0_COLOR    = RGBColor(0xC0, 0x39, 0x2B)     # Muted red
P1_COLOR    = RGBColor(0xD4, 0x8B, 0x2C)     # Muted amber
P2_COLOR    = RGBColor(0x2B, 0x80, 0x5A)     # Muted green

# Theme accent colors (assign one per theme)
THEME_COLORS = [
    RGBColor(0x8B, 0x1A, 0x1A),  # Deep red
    RGBColor(0x1A, 0x5A, 0x8B),  # Deep blue
    RGBColor(0x2B, 0x57, 0x97),  # Medium blue
    RGBColor(0x3A, 0x6B, 0x3A),  # Forest green
    RGBColor(0x7A, 0x4A, 0x1A),  # Warm brown
    RGBColor(0x4A, 0x2B, 0x7A),  # Purple
    RGBColor(0x1A, 0x6B, 0x6B),  # Teal
    RGBColor(0x6B, 0x3A, 0x5A),  # Mauve
    RGBColor(0x4A, 0x6B, 0x2B),  # Olive
]

# ============================================================
# FONTS
# ============================================================
FONT_MAIN  = 'Segoe UI'
FONT_LIGHT = 'Segoe UI Light'
FONT_SEMI  = 'Segoe UI Semibold'

# ============================================================
# SLIDE DIMENSIONS — Widescreen
# ============================================================
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def create_presentation():
    """Create a new widescreen presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs):
    """Add a blank white slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def add_accent_bar(slide, color=ACCENT, width=Inches(0.08)):
    """Add a thin colored bar on the left edge."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bottom_line(slide):
    """Add a subtle bottom separator line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.0), Inches(12.1), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text,
                 font_name=FONT_MAIN, size=Pt(14), color=DARK,
                 bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with a single formatted run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    p.alignment = alignment
    return txBox


def add_hyperlinked_text(slide, left, top, width, height, text, url,
                         font_name=FONT_MAIN, size=Pt(11), color=ACCENT):
    """Add a text box where the text is a clickable hyperlink."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    if url:
        run.hyperlink.address = url
    return txBox


def add_priority_badge(slide, left, top, priority_label, color):
    """Add a small rounded-rectangle priority badge (P0/P1/P2)."""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.5), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = priority_label
    run.font.name = FONT_SEMI
    run.font.size = Pt(9)
    run.font.color.rgb = WHITE
    run.font.bold = True


def get_priority_info(priority_str):
    """Return (label, color) for a priority string."""
    p = str(priority_str)
    if 'P0' in p:
        return 'P0', P0_COLOR
    elif 'P1' in p:
        return 'P1', P1_COLOR
    else:
        return 'P2', P2_COLOR


# ============================================================
# SLIDE TEMPLATES
# ============================================================

def add_title_slide(prs, title, subtitle, detail_line=''):
    """Create the opening title slide."""
    slide = add_blank_slide(prs)
    add_accent_bar(slide)

    add_text_box(slide, Inches(1.0), Inches(2.2), Inches(10), Inches(1.0),
                 title, FONT_LIGHT, Pt(44), DARK)

    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.5), Inches(3), Pt(2)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT
    sep.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(3.9), Inches(10), Inches(0.5),
                 subtitle, FONT_MAIN, Pt(18), GRAY_TEXT)

    if detail_line:
        add_text_box(slide, Inches(1.0), Inches(4.5), Inches(10), Inches(0.5),
                     detail_line, FONT_MAIN, Pt(15), MID_GRAY)

    return slide


def add_theme_divider_slide(prs, theme_name, theme_desc, theme_color, stats_text=''):
    """Create a theme section divider slide."""
    slide = add_blank_slide(prs)
    add_accent_bar(slide, theme_color)

    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1.0),
                 theme_name, FONT_LIGHT, Pt(40), DARK)

    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(2.5), Pt(2)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = theme_color
    sep.line.fill.background()

    if theme_desc:
        add_text_box(slide, Inches(0.8), Inches(4.4), Inches(10), Inches(0.4),
                     theme_desc, FONT_MAIN, Pt(16), GRAY_TEXT)

    return slide


def add_poll_slide(prs, theme_name, theme_color, poll_num, question, options):
    """Create a single polling question slide."""
    slide = add_blank_slide(prs)
    add_accent_bar(slide, theme_color)

    # Theme label
    add_text_box(slide, Inches(0.7), Inches(0.35), Inches(10), Inches(0.35),
                 theme_name, FONT_MAIN, Pt(13), MID_GRAY)

    # Separator
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.8), Inches(11.5), Pt(1)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = LIGHT_GRAY
    sep.line.fill.background()

    # Poll badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.2), Inches(1.1), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = theme_color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'Poll {poll_num}'
    run.font.name = FONT_SEMI
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE
    run.font.bold = True

    # Question text
    add_text_box(slide, Inches(0.7), Inches(1.9), Inches(11.5), Inches(1.2),
                 question, FONT_SEMI, Pt(24), DARK)

    # Answer options
    y = Inches(3.4)
    for i, opt in enumerate(options):
        letter = chr(65 + i)

        # Letter circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.2), y + Pt(2), Inches(0.42), Inches(0.42)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = OFF_WHITE
        circle.line.color.rgb = LIGHT_GRAY
        circle.line.width = Pt(1)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = letter
        run.font.name = FONT_SEMI
        run.font.size = Pt(13)
        run.font.color.rgb = theme_color
        run.font.bold = True

        # Option text
        add_text_box(slide, Inches(1.85), y, Inches(10), Inches(0.5),
                     opt, FONT_MAIN, Pt(16), DARK)

        y += Inches(0.7)

    return slide


# ============================================================
# EXAMPLE USAGE
# ============================================================
if __name__ == '__main__':
    prs = create_presentation()

    # Title slide
    add_title_slide(prs, 'Feedback Analysis', 'Product Area — Feedback Themes',
                    '58 Items • 7 Themes')

    # Example theme divider
    add_theme_divider_slide(prs, 'Security & Access Control',
                            'Network security, data protection, admin controls',
                            THEME_COLORS[0])

    # Example poll slide
    add_poll_slide(prs, 'Security & Access Control', THEME_COLORS[0], 1,
                   'What is your biggest security concern?',
                   ['Lack of vNet support', 'No feature-level controls',
                    'Secret management gaps', 'Sharing beyond security groups'])

    prs.save('example_presentation.pptx')
    print('Saved: example_presentation.pptx')
